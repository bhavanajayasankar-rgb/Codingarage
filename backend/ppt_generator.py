import io
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def format_inr(value: float) -> str:
    return f"\u20B9{value:,.2f}"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

DARK = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x47, 0x55, 0x69)
MUTED = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0x06, 0xB6, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)


def _add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def _add_shape(slide, left, top, width, height, color=LIGHT_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=SLATE, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def _add_title_bar(slide, title: str):
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), ACCENT)
    _add_textbox(slide, Inches(0.6), Inches(0.18), Inches(12), Inches(0.7), title,
                 font_size=28, bold=True, color=WHITE)


def _add_bullets(tf, items, font_size=13, color=SLATE):
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  \u2022  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(4)


def _build_title_slide(prs, report_data, persona):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, ACCENT)

    title = report_data.get("title", f"{persona.capitalize()} Report")
    _add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2), title,
                 font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.6),
                 f"{persona.capitalize()} Persona  |  {datetime.now().strftime('%B %d, %Y')}",
                 font_size=18, color=RGBColor(0xE0, 0xF2, 0xFE), align=PP_ALIGN.CENTER)

    _add_shape(slide, Inches(5.5), Inches(4.5), Inches(2.3), Inches(0.06), WHITE)

    summary = report_data.get("summary", "") or report_data.get("focus", "")
    if summary:
        _add_textbox(slide, Inches(1.5), Inches(4.9), Inches(10), Inches(1.5), summary,
                     font_size=15, color=WHITE, align=PP_ALIGN.CENTER)


def _build_summary_slide(prs, report_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_title_bar(slide, "Executive Summary")

    summary = report_data.get("summary", "") or report_data.get("focus", "")
    cadence = report_data.get("cadence", "N/A")

    _add_textbox(slide, Inches(0.6), Inches(1.4), Inches(4), Inches(0.4), f"Cadence: {cadence}",
                 font_size=13, color=MUTED)

    if summary:
        _add_textbox(slide, Inches(0.6), Inches(2.0), Inches(12), Inches(4.5), summary,
                     font_size=15, color=SLATE)


def _build_metrics_slide(prs, report_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_title_bar(slide, "Key Metrics")

    metrics = report_data.get("metrics", {})
    if not metrics:
        _add_textbox(slide, Inches(0.6), Inches(2.0), Inches(10), Inches(0.5), "No metrics available.",
                     font_size=14, color=MUTED)
        return

    items = list(metrics.items())
    cols = 3
    card_w = Inches(3.8)
    card_h = Inches(1.6)
    start_x = Inches(0.6)
    start_y = Inches(1.6)
    gap_x = Inches(0.3)
    gap_y = Inches(0.3)

    for i, (key, val) in enumerate(items):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        _add_shape(slide, x, y, card_w, card_h, LIGHT_BG)

        label = key.replace("_", " ").title()
        _add_textbox(slide, x + Inches(0.2), y + Inches(0.2), card_w - Inches(0.4), Inches(0.35), label,
                     font_size=11, color=MUTED)

        if isinstance(val, (int, float)):
            if "cost" in key or "savings" in key or "impact" in key or "projected" in key:
                val_str = format_inr(val)
            else:
                val_str = f"{val:,}" if val == int(val) else str(val)
        else:
            val_str = str(val)

        _add_textbox(slide, x + Inches(0.2), y + Inches(0.6), card_w - Inches(0.4), Inches(0.7), val_str,
                     font_size=22, bold=True, color=DARK)


def _build_insights_slide(prs, report_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_title_bar(slide, "Key Insights")

    insights = report_data.get("insights", [])
    if not insights:
        _add_textbox(slide, Inches(0.6), Inches(2.0), Inches(10), Inches(0.5), "No insights available.",
                     font_size=14, color=MUTED)
        return

    txBox = _add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(5.5), "",
                         font_size=14, color=SLATE)
    _add_bullets(txBox.text_frame, insights)


def _build_recommendations_slide(prs, report_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_title_bar(slide, "Recommendations")

    recommendations = report_data.get("recommendations", [])
    if not recommendations:
        _add_textbox(slide, Inches(0.6), Inches(2.0), Inches(10), Inches(0.5), "No recommendations available.",
                     font_size=14, color=MUTED)
        return

    txBox = _add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(5.5), "",
                         font_size=14, color=SLATE)
    _add_bullets(txBox.text_frame, recommendations)


def _build_action_items_slide(prs, report_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_title_bar(slide, "Action Items")

    actions = report_data.get("action_items", [])
    if not actions:
        _add_textbox(slide, Inches(0.6), Inches(2.0), Inches(10), Inches(0.5), "No action items available.",
                     font_size=14, color=MUTED)
        return

    txBox = _add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(5.5), "",
                         font_size=14, color=SLATE)
    _add_bullets(txBox.text_frame, actions)


def _build_alerts_slide(prs, report_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_title_bar(slide, "Alerts")

    alerts = report_data.get("alerts", [])
    if not alerts:
        _add_textbox(slide, Inches(0.6), Inches(2.0), Inches(10), Inches(0.5), "No alerts.",
                     font_size=14, color=MUTED)
        return

    txBox = _add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(5.5), "",
                         font_size=14, color=SLATE)
    alert_texts = []
    for alert in alerts:
        service = alert.get("service", "Unknown")
        msg = alert.get("alert", "")
        alert_texts.append(f"[{service}] {msg}")
    _add_bullets(txBox.text_frame, alert_texts)


def generate_report_ppt(report_data: dict, persona: str) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _build_title_slide(prs, report_data, persona)
    _build_summary_slide(prs, report_data)
    _build_metrics_slide(prs, report_data)
    _build_insights_slide(prs, report_data)
    _build_recommendations_slide(prs, report_data)
    _build_action_items_slide(prs, report_data)
    _build_alerts_slide(prs, report_data)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
